import pandas as pd
from langchain_community.document_loaders import AsyncHtmlLoader
from langchain_community.document_transformers import Html2TextTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tqdm import tqdm
import tiktoken
import matplotlib.pyplot as plt
import logging
from typing import Iterable
from langchain.schema import Document
import os
import io
import requests
import argparse
from pptx import Presentation
from pdfminer.high_level import extract_text
from fake_useragent import UserAgent
from dotenv import load_dotenv
import urllib.parse

#DATA PREPERATION

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)

load_dotenv()

class DataEngine:
    def __init__(self, input_file = "10_links.csv", output_file = "data.jsonl"):
        self.input_file = input_file
        self.output_file = output_file
        self.ua = UserAgent()

    def get_links(self):
        if self.input_file.endswith(".csv"):
            df = pd.read_csv(self.input_file)
            return df['doc_urls'].tolist()
        elif self.input_file.endswith(".txt"):
            with open(self.input_file, "r") as file:
                urls = file.readlines()
                urls = [url.strip() for url in urls if url.strip()]
                return urls

    def filter_urls(self, urls):
        filtered_urls = []
        remaining_urls = []
    
        try :
            for link in urls:
                parsed_url = urllib.parse.urlparse(link)
                path = parsed_url.path
                try:
                    headers = {'User-Agent': self.ua.random}
                    response = requests.get(link, headers=headers, verify=True)
                except requests.exceptions.SSLError:
                    print(f"Ignoring URL due to SSL verification error: {link}")
                    continue
                except requests.exceptions.RequestException as e:
                    print(f"Request error for URL {link}: {e}")
                    continue

                if response.status_code == 200 and not path.endswith((".mp4", ".mp3", ".pdf", ".pptx", ".ppt")):
                    filtered_urls.append(link)
                elif response.status_code == 200 and path.endswith((".pdf", ".pptx", ".ppt")):
                    remaining_urls.append(link)
                
                else:
                    print("inside the else statement")
                    logger.warning(f"Failed to fetch URL: {link}")

            print("Number of filtered urls",len(filtered_urls))
            print("Number of remaining urls",len(remaining_urls))
            
            return filtered_urls, remaining_urls
        
        except Exception as e:
            logger.error(f"Error for some links :{e}")
            return [], []

    def pdf_pptx_text_extractor(self, remaining_urls):
        pdfdocs = []
        pptdocs = []
        pdf_count = 0
        ppt_count = 0
        
        if len(remaining_urls) > 0:
            for url in remaining_urls:
                if url.endswith('.pdf'):
                    pdf_count += 1
                    print("Found PDF URL.................start collecting the data")
                    print("Source :", pdf_count, url)
                    try:
                        response = requests.get(url)
                        if response.status_code == 200:
                            text = extract_text(io.BytesIO(response.content))
                            pdfdocs.append(Document(page_content=text, metadata={'source': url}))
                        else:
                            logger.warning(f"Failed to fetch URL: {url}")
                    except Exception as e:
                        logger.error(f"Error extracting text from PDF: {e}")

                elif url.endswith('.pptx'):
                    ppt_count += 1
                    print("Found PPT URL.................start collecting the data")
                    print("Source :", ppt_count, url)
                    try:
                        response = requests.get(url)
                        if response.status_code == 200:
                            pptx_file = io.BytesIO(response.content)
                            presentation = Presentation(pptx_file)
                            text = ""
                            for slide in presentation.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            pptdocs.append(Document(page_content=text, metadata={'source': url}))
                        else:
                            logger.warning(f"Failed to fetch URL: {url}")
                    except Exception as e:
                        logger.error(f"Error extracting text from PowerPoint file: {e}")
        
        return pdfdocs, pptdocs


    def web_text_extractor(self, urls):
        filtered_urls, remaining_urls = self.filter_urls(urls)
        docs_transformed = None
        
        try:
            urls = filtered_urls
            logger.info("Extracting raw data (html) from URLs...")
            loader = AsyncHtmlLoader(urls, raise_for_status=True, requests_per_second=100)
            docs = loader.load()
            logger.info(f"Raw documents got from the collected URLs are : {len(docs)}")
            html2text = Html2TextTransformer()
            docs_transformed = html2text.transform_documents(docs)
            logger.info("Tranforming the documents from html to text is DONE...")
            pdftext, ppttext = self.pdf_pptx_text_extractor(remaining_urls)
            document_all = docs_transformed + pdftext + ppttext
            return document_all

        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            return None

    def _text_splitter(self, docs):
        text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name="cl100k_base",
            chunk_size=5000,
            chunk_overlap=500
        )
        doc_splits = text_splitter.split_documents(docs)
        logger.info("Data chunking done successfully.")
        return doc_splits
    

    def num_tokens_from_string(self, string: str, encoding_name: str) -> int:
        encoding = tiktoken.get_encoding(encoding_name)
        num_tokens = len(encoding.encode(string))
        return num_tokens

    def save_docs_to_jsonl(self, array:Iterable[Document])->None:
        with open(self.output_file, 'w') as jsonl_file:
            for doc in tqdm(array, desc="saving data :"):
                jsonl_file.write(doc.json() + '\n')

        logger.info(f"Data saved to {self.output_file}")

    def append_to_file(self, file_path: str, content: str) -> None:
        with open(file_path, 'w') as file:
            file.write(content)
        print("Data saved succesfully--------------------")

    @staticmethod
    def _plot_str_count(str_counts, path):
        plt.figure(figsize=(10, 6))
        plt.hist(str_counts, bins=30, color="blue", edgecolor="black", alpha=0.7)
        plt.title("Histogram of Token Counts")
        plt.xlabel("Token Count")
        plt.ylabel("Frequency")
        plt.grid(axis="y", alpha=0.75)
        plt.savefig(path)

if __name__ == "__main__":

    parser = argparse.ArgumentParser(description="Web Scraper")
    parser.add_argument("--input_file", type=str, default="Extracted_links.csv", help="Path for the links")
    parser.add_argument("--output_file", type=str, default="docs.jsonl", help="Path to save the extracted data")
    args = parser.parse_args()

    data_engine = DataEngine(args.input_file, args.output_file)
    urls = data_engine.get_links()
    docs = data_engine.web_text_extractor(urls)
    data_engine.save_docs_to_jsonl(docs)

 
'''
Usage :

--> provide path for inputfile of urls and path for saving data to a csv file.
--> provide a filepath for plotting the data distribution

python3 data_processing.py

'''