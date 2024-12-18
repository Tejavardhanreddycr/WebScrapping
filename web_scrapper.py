"""
Web Scraping Data Engine

This module provides functionality for scraping and processing web content from various sources
including HTML pages, PDFs, and PowerPoint presentations. It handles data extraction,
transformation, and storage in JSONL format.

"""

from typing import Iterable, List, Tuple, Optional
import argparse
import io
import logging
import os
import urllib.parse

# Third-party imports
import pandas as pd
import requests
import tiktoken
import matplotlib.pyplot as plt
from dotenv import load_dotenv
from fake_useragent import UserAgent
from langchain.schema import Document
from langchain_community.document_loaders import AsyncHtmlLoader
from langchain_community.document_transformers import Html2TextTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pdfminer.high_level import extract_text
from pptx import Presentation
from tqdm import tqdm

# Constants
DEFAULT_CHUNK_SIZE = 5000
DEFAULT_CHUNK_OVERLAP = 500
DEFAULT_ENCODING = "cl100k_base"
SUPPORTED_INPUT_FORMATS = ('.csv', '.txt')
SUPPORTED_DOC_FORMATS = ('.pdf', '.pptx', '.ppt')

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

class DataEngine:
    """A class for extracting and processing web content from various sources.

    This class handles the extraction of text from web pages, PDFs, and PowerPoint
    presentations, processes the extracted text, and saves it in a structured format.

    Attributes:
        input_file (str): Path to the input file containing URLs
        output_file (str): Path to save the processed data
        ua (UserAgent): User agent for making HTTP requests
    """

    def __init__(self, input_file: str = "10_links.csv", output_file: str = "data.jsonl"):
        """Initialize the DataEngine.

        Args:
            input_file: Path to the file containing URLs (CSV or TXT)
            output_file: Path where the processed data will be saved
        """
        self.input_file = input_file
        self.output_file = output_file
        self.ua = UserAgent()

    def get_links(self) -> List[str]:
        """Extract URLs from the input file.

        Returns:
            List of URLs extracted from the input file.

        Raises:
            ValueError: If the input file format is not supported
        """
        file_ext = os.path.splitext(self.input_file)[1].lower()
        if file_ext not in SUPPORTED_INPUT_FORMATS:
            raise ValueError(f"Unsupported input file format. Supported formats: {SUPPORTED_INPUT_FORMATS}")

        if file_ext == '.csv':
            df = pd.read_csv(self.input_file)
            return df['doc_urls'].tolist()
        else:  # .txt file
            with open(self.input_file, "r") as file:
                return [url.strip() for url in file.readlines() if url.strip()]

    def filter_urls(self, urls: List[str]) -> Tuple[List[str], List[str]]:
        """Filter URLs based on their accessibility and file type.

        Args:
            urls: List of URLs to filter

        Returns:
            Tuple containing filtered URLs and remaining URLs
        """
        filtered_urls = []
        remaining_urls = []
    
        for link in urls:
            try:
                parsed_url = urllib.parse.urlparse(link)
                path = parsed_url.path.lower()
                headers = {'User-Agent': self.ua.random}
                
                try:
                    response = requests.get(link, headers=headers, verify=True, timeout=30)
                    response.raise_for_status()
                except requests.exceptions.SSLError:
                    logger.warning(f"SSL verification error: {link}")
                    continue
                except requests.exceptions.RequestException as e:
                    logger.warning(f"Request error for URL {link}: {e}")
                    continue

                if path.endswith(SUPPORTED_DOC_FORMATS):
                    remaining_urls.append(link)
                else:
                    filtered_urls.append(link)
                
            except Exception as e:
                logger.error(f"Error processing URL {link}: {e}")

        logger.info(f"Filtered URLs: {len(filtered_urls)}, Remaining URLs: {len(remaining_urls)}")
        return filtered_urls, remaining_urls

    def pdf_pptx_text_extractor(self, remaining_urls: List[str]) -> Tuple[List[Document], List[Document]]:
        """Extract text from PDF and PowerPoint files.

        Args:
            remaining_urls: List of URLs pointing to PDF or PowerPoint files

        Returns:
            Tuple containing lists of processed PDF and PowerPoint documents
        """
        pdfdocs = []
        pptdocs = []
        
        for url in remaining_urls:
            try:
                response = requests.get(url, timeout=30)
                response.raise_for_status()
                
                if url.lower().endswith('.pdf'):
                    logger.info(f"Processing PDF: {url}")
                    text = extract_text(io.BytesIO(response.content))
                    pdfdocs.append(Document(page_content=text, metadata={'source': url}))
                
                elif url.lower().endswith(('.pptx', '.ppt')):
                    logger.info(f"Processing PowerPoint: {url}")
                    pptx_file = io.BytesIO(response.content)
                    presentation = Presentation(pptx_file)
                    text = "\n".join(
                        shape.text
                        for slide in presentation.slides
                        for shape in slide.shapes
                        if hasattr(shape, "text")
                    )
                    pptdocs.append(Document(page_content=text, metadata={'source': url}))
                
            except Exception as e:
                logger.error(f"Error processing document {url}: {e}")

        return pdfdocs, pptdocs

    def web_text_extractor(self, urls: List[str]) -> Optional[List[Document]]:
        """Extract text content from web pages and documents.

        Args:
            urls: List of URLs to process

        Returns:
            List of processed documents or None if extraction fails
        """
        filtered_urls, remaining_urls = self.filter_urls(urls)
        
        try:
            logger.info("Extracting HTML content...")
            loader = AsyncHtmlLoader(filtered_urls, raise_for_status=True, requests_per_second=100)
            docs = loader.load()
            logger.info(f"Processed {len(docs)} HTML documents")

            html2text = Html2TextTransformer()
            docs_transformed = html2text.transform_documents(docs)
            
            pdftext, ppttext = self.pdf_pptx_text_extractor(remaining_urls)
            return docs_transformed + pdftext + ppttext

        except Exception as e:
            logger.error(f"Error in web text extraction: {e}")
            return None

    def _text_splitter(self, docs: List[Document]) -> List[Document]:
        """Split documents into smaller chunks for processing.

        Args:
            docs: List of documents to split

        Returns:
            List of split documents
        """
        text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name=DEFAULT_ENCODING,
            chunk_size=DEFAULT_CHUNK_SIZE,
            chunk_overlap=DEFAULT_CHUNK_OVERLAP
        )
        return text_splitter.split_documents(docs)

    @staticmethod
    def num_tokens_from_string(string: str, encoding_name: str = DEFAULT_ENCODING) -> int:
        """Count the number of tokens in a string.

        Args:
            string: Input text
            encoding_name: Name of the tokenizer encoding

        Returns:
            Number of tokens in the string
        """
        encoding = tiktoken.get_encoding(encoding_name)
        return len(encoding.encode(string))

    def save_docs_to_jsonl(self, documents: Iterable[Document]) -> None:
        """Save documents to a JSONL file.

        Args:
            documents: Documents to save
        """
        if not documents:
            logger.warning("No documents to save")
            return

        with open(self.output_file, 'w') as jsonl_file:
            for doc in tqdm(documents, desc="Saving documents"):
                jsonl_file.write(doc.json() + '\n')

        logger.info(f"Documents saved to {self.output_file}")

    @staticmethod
    def _plot_str_count(str_counts: List[int], path: str) -> None:
        """Generate and save a histogram of token counts.

        Args:
            str_counts: List of token counts
            path: Path to save the plot
        """
        plt.figure(figsize=(10, 6))
        plt.hist(str_counts, bins=30, color="blue", edgecolor="black", alpha=0.7)
        plt.title("Histogram of Token Counts")
        plt.xlabel("Token Count")
        plt.ylabel("Frequency")
        plt.grid(axis="y", alpha=0.75)
        plt.savefig(path)
        plt.close()

def main():
    """Main entry point for the web scraper."""
    parser = argparse.ArgumentParser(description="Web content scraper and processor")
    parser.add_argument(
        "--input_file",
        type=str,
        default="Extracted_links.csv",
        help="Path to the file containing URLs (CSV or TXT)"
    )
    parser.add_argument(
        "--output_file",
        type=str,
        default="docs.jsonl",
        help="Path to save the extracted data (JSONL format)"
    )
    args = parser.parse_args()

    data_engine = DataEngine(args.input_file, args.output_file)
    
    try:
        urls = data_engine.get_links()
        if not urls:
            logger.error("No URLs found in the input file")
            return

        docs = data_engine.web_text_extractor(urls)
        if docs:
            data_engine.save_docs_to_jsonl(docs)
        else:
            logger.error("No documents were processed successfully")

    except Exception as e:
        logger.error(f"An error occurred: {e}")

if __name__ == "__main__":
    main()